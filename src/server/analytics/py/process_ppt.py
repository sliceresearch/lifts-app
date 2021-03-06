#!/usr/bin/env python
# coding: utf-8

# # Interactive Powerpoint Analytics Tool (IPAT)
#
# A product of the **LIFTS** project (**Learning in Future Teaching Spaces**), 2019-2020.
#
# Authors:
# * Mark Utting
# * Jacqui Blake

# In[89]:


from pptx import Presentation
import collections
import json
import sys
from typing import List, Dict, Tuple, Any


# In[63]:

MAX_WORDS_PER_SLIDE = 65

# pres_name = "ICT112_Week09_HTML.pptx"
# pres_name = "presentation-template_teaching-lecture_wide_2018.pptx"
pres_name = "pres_test.pptx"


# In[65]:

def debug_dump(prs:Presentation):
    """Print lots of basic information about the given presentation."""
    print("Presentation has", len(prs.slides), "slides")

    # Print summary of all slides, plus text
    n = 0
    for slide in prs.slides:
        n += 1
        print("========== slide {} ========== [{}]".format(n, slide.slide_layout.name))
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            print(shape.name)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print("    " + run.text)


# %%

def get_word_counts(slides) -> List[int]:
    """Count the amount of text in each slide."""
    word_count = []
    for slide in slides:
        # print(f"========== slide {len(text_count)+1} ========== [{slide.slide_layout.name}]")
        words = 0
        # find all text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # print(shape.name)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # print("    " + run.text)
                    words += len(run.text.split())
        word_count.append(words)
    return word_count


# In[91]:


def calculate_text_stars(word_counts) -> int:
    """Calculates a one-to-five star ranking for presentations that are not too text-heavy."""
    if word_counts == []:
        return 3
    words_per_slide = sum(word_counts) / len(word_counts)
    stars = 5 - abs(words_per_slide - 35) / 8
    # print(stars)
    return max(0, min(5, int(stars + 0.5)))

assert calculate_text_stars([]) == 3
assert calculate_text_stars([30, 40]) == 5
assert calculate_text_stars([10, 10]) == 2
assert calculate_text_stars([10, 30, 60]) == 5  # TODO: penalise text-heavy slides more?


# In[92]:


INTERACTIVE = set(["MS Forms Quiz/Survey", "Video", "Demo"])

def count_layouts(prs:Presentation) -> Tuple[int, Dict[str, int]]:
    """Counts how many times each PPT layout is used.
    Returns the total number of interactive layouts, plus a dictionary of the layout counts.
    """
    layouts = collections.defaultdict(int)
    layouts_interactive = 0
    for slide in prs.slides:
        layouts[slide.slide_layout.name] += 1
        if slide.slide_layout.name in INTERACTIVE:
            layouts_interactive += 1
    return (layouts_interactive, layouts)

##################################################################################presentation properties

def get_presentation_properties(prs:Presentation) -> Dict[str, Any]:
    pres_properties = prs.core_properties
    result = {
        "author":pres_properties.author,
        "title":pres_properties.title,
        "category":pres_properties.category,
        "subject":"00/00/00", #pres_properties.subject,
        "created":"00/00/00", #pres_properties.created,
        "modified":"00/00/00", #pres_properties.modified,
        "last_modified_by":pres_properties.last_modified_by
        }

    return result

##################################################################################slide analytics
def get_slide_analytics(slides):
    slides_out = []
    for slide in slides:
        words = 0
        shapes_out = []
        for shape in slide.shapes:
            shape_analyse = analyse_shape(shape)
            shapes_out.append(shape_analyse);

        result = {
             "id":slide.slide_id,
             "name":slide.name,
             "shapes":shapes_out,
             "words":words,
        }
        slides_out.append(result)
     
    return slides_out

def analyse_shape(shape):
#text
    shape_text=""
    shape_paragraphs=[]
    shape_paragraphs_properties=[]
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:  ## extract text 
           shape_font_properties = analyse_shape_paragraph_properties(paragraph)
           shape_paragraphs.append(paragraph.text)
           shape_paragraphs_properties.append(shape_font_properties)
        shape_text=shape.text  ## all text == all paragraph text?

    shape_type = analyse_shape_type(shape.name);
    shapes_out = {'name':shape.name,'type':shape_type,'paragraphs':shape_paragraphs,'paragraphs_properties':shape_paragraphs_properties}

    return shapes_out

def analyse_shape_paragraph_properties(paragraph):
    shape_level=paragraph.level
    shape_font=paragraph.font
    font_result = {
             "name":shape_font.name,
             "size":shape_font.size
        }
    result = {
             "font":font_result,
             "level":shape_level
        }
    return result

def analyse_shape_type(shape_name):

   # if (shape_type==14):
   #      return 'title'
  #  elif (shape_type==17):
  #       return 'content'
    
    if (shape_name.find('Title') != -1): 
         return 'title'
    elif (shape_name.find('Content') != -1):
         return 'content'

    return 'none'



# In[94]:

def analyse_presentation(pres_name:str, verbose=False) -> Dict[str, Any]:
    """Analyses a presentation and returns a dictionary of star-ratings, plus extra details.
    Works best on presentations that use the LIFTS template, with known layout names.
    """
    prs = Presentation(pres_name)
    if verbose:
        debug_dump(prs)
    (layouts_interactive, layouts) = count_layouts(prs)
    interaction_stars = min(layouts_interactive, 5)
    topic_stars = ([1,1,3,5,5,4,3,2,1]+[1]*100)[layouts["Section Header"]]

    word_count = get_word_counts(prs.slides)
    words_per_slide = sum(word_count) / len(word_count)
    # ideal words/slide is 30-40 (5 stars)
    text_stars = calculate_text_stars(word_count)
    # print("word counts:", word_count)

    # Create a list of warnings about very text-heavy slides
    heavy_warnings = []
    for slide, words in enumerate(word_count):
        if words > MAX_WORDS_PER_SLIDE:
            heavy_warnings.append(f"WARNING: slide {slide} has {words} words!")

    pres_properties = get_presentation_properties(prs)
    slides = get_slide_analytics(prs.slides)

    result = {

        "presentation_rating_stars_interaction": interaction_stars,
        "presentation_rating_stars_section": topic_stars,
        "presentation_rating_stars_accessibility": 3,   # not implemented yet!
        "presentation_rating_stars_text": text_stars,
        "presentation_count_slide": len(prs.slides),
        "presentation_count_layout": layouts,  # dictionary that maps layout name to count
        "presentation_total_words": words_per_slide,  # a float
        "presentation_warning_text_heavy": heavy_warnings,  # a list of warning strings
        "presentation_data_slides": slides,  # a list of slides and analytics
        "filename": pres_name,  # TODO: strip any Path and just return file name?
        "name": "USC999",
   #     "description": "Introduction to ICT",
        "properties": pres_properties
       # "slides": slides,
        }

    return result


# %%
#print(sys.argv[1])
analytics = analyse_presentation(sys.argv[1])  # the default example
jout = json.dumps(analytics)
print(jout)
#if __name__ == "__main__":
#    verbose = False
#    if len(sys.argv) > 1:
#        for arg in sys.argv:
#            if arg == "--verbose":
#                verbose = True
#            else:
#                analytics = analyse_presentation(arg, verbose=verbose)
#                print(json.dumps(analytics, indent=2))
#    else:
#        print("Usage: [--verbose] file1.pptx file2.pptx ...")
#        analytics = analyse_presentation(pres_name)  # the default example
#        print(json.dumps(analytics, indent=2))

